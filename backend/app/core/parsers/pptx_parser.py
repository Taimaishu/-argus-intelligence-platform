"""PPTX presentation parser using python-pptx."""

from pptx import Presentation
from .base import BaseParser, ParsedDocument


class PPTXParser(BaseParser):
    """Parser for PPTX presentations."""

    def supports_file_type(self, file_extension: str) -> bool:
        """Check if file extension is .pptx"""
        return file_extension.lower() in [".pptx", ".ppt"]

    def parse(self, file_path: str) -> ParsedDocument:
        """
        Parse PPTX presentation and extract text.

        Args:
            file_path: Path to PPTX file

        Returns:
            ParsedDocument with extracted content
        """
        try:
            prs = Presentation(file_path)

            # Extract metadata
            core_props = prs.core_properties
            title = core_props.title or "Presentation"
            author = core_props.author

            full_text = ""
            sections = []

            # Process each slide
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = f"Slide {slide_num}\n\n"

                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"

                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(cell.text for cell in row.cells)
                            slide_text += row_text + "\n"

                if slide_text.strip():
                    sections.append(
                        {
                            "heading": f"Slide {slide_num}",
                            "content": slide_text,
                            "page": slide_num,
                        }
                    )

                    full_text += slide_text + "\n\n"

            return ParsedDocument(
                text=self.clean_text(full_text),
                title=title,
                author=author,
                metadata={
                    "slide_count": len(prs.slides),
                    "created": str(core_props.created) if core_props.created else None,
                    "modified": (
                        str(core_props.modified) if core_props.modified else None
                    ),
                    "subject": core_props.subject,
                },
                sections=sections,
                page_count=len(prs.slides),
            )

        except Exception as e:
            raise Exception(f"Failed to parse PPTX: {str(e)}")
